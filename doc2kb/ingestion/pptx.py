"""Extract text from a PPTX presentation slide by slide."""
from __future__ import annotations

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

from ..utils import content_doc_id


def ingest_pptx(path: Path) -> tuple[str, dict]:
    prs = Presentation(str(path))
    lines: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        lines.append(f"## Slide {slide_num}")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                # Treat large text as a heading
                if para.runs and any(
                    run.font.size and run.font.size >= Pt(20) for run in para.runs
                ):
                    lines.append(f"### {text}")
                else:
                    lines.append(text)
        lines.append("")

    md_text = "\n".join(lines)
    doc_id = content_doc_id(path)

    metadata = {
        "title": path.stem,
        "source": str(path.resolve()),
        "type": "pptx",
        "date_ingested": datetime.now().isoformat(timespec="seconds"),
        "tags": [],
        "doc_id": doc_id,
    }
    return md_text, metadata
